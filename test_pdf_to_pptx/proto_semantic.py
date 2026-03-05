"""Prototype semantic: PDF → PPTX via element extraction (text, images, shapes)."""

import io
import argparse
import os
import fitz
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches
from pptx.oxml.ns import qn
from pptx.enum.dml import MSO_THEME_COLOR
import lxml.etree as etree


def int_to_rgb(color_int: int) -> RGBColor:
    """Convert fitz integer color (0xRRGGBB) to RGBColor."""
    r = (color_int >> 16) & 0xFF
    g = (color_int >> 8) & 0xFF
    b = color_int & 0xFF
    return RGBColor(r, g, b)


def fitz_color_to_rgb(color) -> RGBColor:
    """Handle fitz color: can be float tuple (r,g,b) or int."""
    if color is None:
        return RGBColor(0, 0, 0)
    if isinstance(color, (list, tuple)):
        if len(color) == 3:
            r, g, b = [int(c * 255) for c in color]
            return RGBColor(r, g, b)
        if len(color) == 1:
            v = int(color[0] * 255)
            return RGBColor(v, v, v)
    if isinstance(color, (int, float)):
        v = int(color)
        return int_to_rgb(v)
    return RGBColor(0, 0, 0)


def get_slide_background_color(page):
    """Return RGBColor of largest filled rect covering >80% of page, or None."""
    page_area = page.rect.width * page.rect.height
    best_fill = None
    best_area = 0
    for path in page.get_drawings():
        fill = path.get("fill")
        rect = path.get("rect")
        if fill is None or rect is None:
            continue
        if fill in [(1.0, 1.0, 1.0), (0.0, 0.0, 0.0), None]:
            continue
        area = rect.width * rect.height
        if area > 0.8 * page_area and area > best_area:
            best_area = area
            best_fill = fill
    if best_fill is None:
        return None
    return fitz_color_to_rgb(best_fill)


def add_drawings(slide, page, scale_x, scale_y, bg_color=None):
    """Add filled rectangles from vector drawings as background shapes."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    drawings = page.get_drawings()
    rects_added = 0
    for path in drawings:
        fill = path.get("fill")
        rect = path.get("rect")
        if fill is None or rect is None:
            continue
        # Always skip pure black border artifacts
        if fill == (0.0, 0.0, 0.0):
            continue
        # Skip white only when slide background is also white (no colored bg)
        is_white = fill == (1.0, 1.0, 1.0)
        if is_white and bg_color is None:
            continue
        # Skip rects that bleed outside page bounds
        pw, ph = page.rect.width, page.rect.height
        if rect.x0 < -5 or rect.y0 < -5 or rect.x1 > pw + 5 or rect.y1 > ph + 5:
            continue
        # Skip micro-surfaces
        if rect.width * rect.height < 500:
            continue
        left = Pt(rect.x0 * scale_x)
        top = Pt(rect.y0 * scale_y)
        width = Pt(rect.width * scale_x)
        height = Pt(rect.height * scale_y)
        shape = slide.shapes.add_shape(
            1, left, top, width, height  # MSO_SHAPE_TYPE.RECTANGLE = 1
        )
        fill_color = fitz_color_to_rgb(fill)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()  # no border
        rects_added += 1

    # Move rect shapes to back (before text shapes)
    sp_tree = slide.shapes._spTree
    all_shapes = list(sp_tree)
    # shapes added are at the end; move the last `rects_added` before other shapes
    # Actually we want rects behind text, so insert them early in spTree
    # The first 2 children are nvGrpSpPr and grpSpPr — keep those first
    if rects_added > 0:
        non_group = [c for c in all_shapes if c.tag.endswith("}sp")]
        rect_els = non_group[-rects_added:]
        for el in rect_els:
            sp_tree.remove(el)
            sp_tree.insert(2, el)


def add_images(slide, page, doc, scale_x, scale_y):
    """Add embedded images to slide."""
    images = page.get_images(full=True)
    for img_info in images:
        xref = img_info[0]
        try:
            rects = page.get_image_rects(xref)
            if not rects:
                continue
            rect = rects[0]
            # Clamp to page bounds
            clipped = fitz.Rect(
                max(rect.x0, 0), max(rect.y0, 0),
                min(rect.x1, page.rect.width), min(rect.y1, page.rect.height)
            )
            if clipped.is_empty or clipped.width < 5 or clipped.height < 5:
                continue
            img_data = doc.extract_image(xref)
            if not img_data:
                continue
            img_bytes = img_data["image"]
            left = Pt(clipped.x0 * scale_x)
            top = Pt(clipped.y0 * scale_y)
            width = Pt(clipped.width * scale_x)
            height = Pt(clipped.height * scale_y)
            slide.shapes.add_picture(io.BytesIO(img_bytes), left, top, width, height)
        except Exception as e:
            print(f"  [warn] image xref={xref}: {e}")


FONT_MAP = {
    "liberationsans":  "Arial",
    "liberationserif": "Times New Roman",
    "liberationmono":  "Courier New",
    "dejavusans":      "Arial",
    "dejavuserif":     "Times New Roman",
    "dejavumono":      "Courier New",
    "helvetica":       "Arial",
    "times":           "Times New Roman",
    "courier":         "Courier New",
    "inter":           "Liberation Sans",
    "inter24pt":       "Liberation Sans",
    "inter28pt":       "Liberation Sans",
}


def map_font_name(pdf_font: str) -> str:
    """Map PDF font name to a common Windows/PowerPoint font."""
    key = pdf_font.lower().split("-")[0].replace(" ", "")
    return FONT_MAP.get(key, "Arial")  # default to Arial


def dominant_font_size(block) -> float:
    sizes = [s["size"] for l in block.get("lines", []) for s in l.get("spans", [])]
    return sum(sizes) / len(sizes) if sizes else 12


def merge_text_blocks(blocks: list) -> list:
    """Merge vertically adjacent text blocks that belong to the same paragraph."""
    import copy
    if not blocks:
        return blocks

    # Sort by column (x0 bucketed to 40pt) then y0
    sorted_blocks = sorted(blocks, key=lambda b: (b["bbox"][0] // 40, b["bbox"][1]))
    merged = []
    for b in sorted_blocks:
        if not merged:
            merged.append(copy.deepcopy(b))
            continue
        prev = merged[-1]
        px0, py0, px1, py1 = prev["bbox"]
        bx0, by0, bx1, by1 = b["bbox"]
        same_col = abs(px0 - bx0) < 30
        small_gap = (by0 - py1) < 15
        same_size = abs(dominant_font_size(prev) - dominant_font_size(b)) <= 1
        if same_col and small_gap and same_size:
            # Extend bbox and append lines
            prev["bbox"] = (min(px0, bx0), py0, max(px1, bx1), by1)
            prev["lines"].extend(copy.deepcopy(b.get("lines", [])))
        else:
            merged.append(copy.deepcopy(b))
    return merged


def add_text_blocks(slide, page, scale_x, scale_y):
    """Add text blocks as textboxes, grouping spans by block bbox."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    raw_blocks = [b for b in page.get_text("dict")["blocks"] if b.get("type") == 0]
    blocks = merge_text_blocks(raw_blocks)
    for block in blocks:
        bbox = block["bbox"]  # (x0, y0, x1, y1)
        left = Pt(bbox[0] * scale_x)
        top = Pt(bbox[1] * scale_y)
        width = Pt((bbox[2] - bbox[0]) * scale_x)
        height = Pt((bbox[3] - bbox[1]) * scale_y)

        # Ensure minimum dimensions
        if width < Pt(1) or height < Pt(1):
            continue

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = False
        # Enable spAutoFit so box expands to fit text without clipping
        body_pr = tf._txBody.bodyPr
        etree.SubElement(body_pr, qn('a:spAutoFit'))
        for el in body_pr.findall(qn('a:noAutofit')):
            body_pr.remove(el)

        first_para = True
        for line in block.get("lines", []):
            if first_para:
                para = tf.paragraphs[0]
                first_para = False
            else:
                para = tf.add_paragraph()

            for span in line.get("spans", []):
                text = span.get("text", "")
                if not text.strip():
                    continue
                run = para.add_run()
                run.text = text
                font_name = span.get("font", "")
                size = span.get("size", 12)
                color_val = span.get("color", 0)

                run.font.size = Pt(size)
                fn_lower = font_name.lower()
                run.font.name = map_font_name(font_name)
                run.font.bold = "bold" in fn_lower
                run.font.italic = "italic" in fn_lower or "oblique" in fn_lower

                # fitz color: integer 0xRRGGBB
                if isinstance(color_val, (int, float)):
                    run.font.color.rgb = int_to_rgb(int(color_val))
                elif isinstance(color_val, (list, tuple)):
                    run.font.color.rgb = fitz_color_to_rgb(color_val)


def convert(pdf_path: str, output_path: str | None = None) -> str:
    """Convert a PDF file to PPTX. Returns output path."""
    if output_path is None:
        base = os.path.splitext(pdf_path)[0]
        output_path = base + "_semantic.pptx"

    doc = fitz.open(pdf_path)
    prs = Presentation()

    # Use first page dimensions to set slide size
    first_page = doc[0]
    pw = first_page.rect.width
    ph = first_page.rect.height
    prs.slide_width = Pt(pw)
    prs.slide_height = Pt(ph)

    blank_layout = prs.slide_layouts[6]  # blank layout

    for page_num, page in enumerate(doc):
        slide = prs.slides.add_slide(blank_layout)

        # Scale factors (PDF pts → PPTX pts, usually 1:1 but handle size diff)
        scale_x = pw / page.rect.width
        scale_y = ph / page.rect.height

        # Order: drawings (background) → images → text (foreground)
        bg_color = get_slide_background_color(page)
        if bg_color:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = bg_color
        add_drawings(slide, page, scale_x, scale_y, bg_color=bg_color)
        add_images(slide, page, doc, scale_x, scale_y)
        add_text_blocks(slide, page, scale_x, scale_y)

        shapes = slide.shapes
        print(f"\nSlide {page_num + 1}: {len(shapes)} shapes found")
        for i, shape in enumerate(shapes):
            stype = shape.shape_type
            type_name = {1: "AUTO_SHAPE", 13: "PICTURE", 17: "TEXT_BOX"}.get(
                stype, f"TYPE_{stype}"
            )
            if stype == 17:  # text box
                text_preview = shape.text_frame.text[:40].replace("\n", " ")
                print(f'  - shape {i}: type={type_name}, text="{text_preview}"')
            else:
                print(f"  - shape {i}: type={type_name}")

    prs.save(output_path)
    doc.close()
    print(f"\nSaved: {output_path}")
    return output_path


if __name__ == "__main__":
    pdf_path = os.path.join(os.path.dirname(__file__), "test_pdf_ia.pdf")
    pptx_path = os.path.join(
        os.path.dirname(__file__), "output", "test_output_semantic.pptx"
    )

    out = convert(pdf_path, pptx_path)
    print(f"PPTX saved: {out}")
    print(f"File size: {os.path.getsize(pptx_path):,} bytes")

"""Tests for the hybrid PDF → PPTX converter (proto_hybrid.py)."""
import os
import sys
import tempfile
import pytest
import fitz
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Allow importing from parent dir
sys.path.insert(0, os.path.dirname(os.path.dirname(__file__)))

from proto_hybrid import convert

PDF_PATH = os.path.join(os.path.dirname(os.path.dirname(__file__)), "test_pdf_ia.pdf")
SSIM_THRESHOLD = 0.87  # minimum acceptable mean SSIM (text erased from bg reduces SSIM slightly)


@pytest.fixture(scope="module")
def converted_pptx(tmp_path_factory):
    out = tmp_path_factory.mktemp("output") / "test_hybrid_out.pptx"
    convert(PDF_PATH, str(out))
    return str(out)


def test_output_file_exists(converted_pptx):
    assert os.path.isfile(converted_pptx)
    assert os.path.getsize(converted_pptx) > 10_000  # at least 10 KB


def test_slide_count_matches_pdf(converted_pptx):
    doc = fitz.open(PDF_PATH)
    pdf_pages = len(doc)
    doc.close()
    prs = Presentation(converted_pptx)
    assert len(prs.slides) == pdf_pages


def test_slide_dimensions_match_pdf(converted_pptx):
    doc = fitz.open(PDF_PATH)
    page = doc[0]
    expected_w_emu = int(page.rect.width * 12700)   # Pt(n) = n * 12700 EMU
    expected_h_emu = int(page.rect.height * 12700)
    doc.close()
    prs = Presentation(converted_pptx)
    assert abs(prs.slide_width - expected_w_emu) < 50
    assert abs(prs.slide_height - expected_h_emu) < 50


def test_first_shape_is_background_image(converted_pptx):
    """Each slide must start with a full-page PICTURE at position (0,0)."""
    prs = Presentation(converted_pptx)
    for i, slide in enumerate(prs.slides):
        first = slide.shapes[0]
        assert first.shape_type == MSO_SHAPE_TYPE.PICTURE, (
            f"Slide {i+1}: first shape is {first.shape_type}, expected PICTURE"
        )
        assert first.left == 0
        assert first.top == 0
        assert first.width == prs.slide_width
        assert first.height == prs.slide_height


def test_has_text_shapes(converted_pptx):
    """At least one TEXT_BOX per slide (unless slide has no text at all)."""
    prs = Presentation(converted_pptx)
    for i, slide in enumerate(prs.slides):
        text_shapes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX]
        assert len(text_shapes) >= 1, f"Slide {i+1} has no text boxes"


def test_text_boxes_are_transparent(converted_pptx):
    """Text boxes must not have a solid fill (transparent so background shows)."""
    from pptx.enum.dml import MSO_THEME_COLOR
    from pptx.oxml.ns import qn
    prs = Presentation(converted_pptx)
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                # Check there's no solid fill element in the spPr
                sp = shape._element
                spPr = sp.find(qn("p:spPr"))
                if spPr is not None:
                    solidFill = spPr.find(f".//{qn('a:solidFill')}")
                    assert solidFill is None, (
                        f"Slide {i+1} text box has solid fill (should be transparent)"
                    )


def test_ssim_above_threshold(converted_pptx):
    """Mean SSIM between source PDF and LibreOffice-rendered PPTX must be above threshold."""
    import subprocess
    import numpy as np
    from skimage.metrics import structural_similarity as ssim
    from skimage import color

    # Render PPTX → PDF via LibreOffice
    out_dir = os.path.dirname(converted_pptx)
    result = subprocess.run(
        ["libreoffice", "--headless", "--convert-to", "pdf",
         "--outdir", out_dir, converted_pptx],
        capture_output=True, text=True
    )
    pptx_name = os.path.splitext(os.path.basename(converted_pptx))[0]
    pptx_pdf = os.path.join(out_dir, pptx_name + ".pdf")
    assert os.path.isfile(pptx_pdf), f"LibreOffice failed: {result.stderr}"

    src_doc  = fitz.open(PDF_PATH)
    conv_doc = fitz.open(pptx_pdf)
    assert len(src_doc) == len(conv_doc)

    scores = []
    for page_num in range(len(src_doc)):
        def render(doc, n):
            pix = doc[n].get_pixmap(matrix=fitz.Matrix(1.5, 1.5))
            arr = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, 3)
            return color.rgb2gray(arr)

        src_img  = render(src_doc, page_num)
        conv_img = render(conv_doc, page_num)
        # Resize if necessary
        if src_img.shape != conv_img.shape:
            from skimage.transform import resize
            conv_img = resize(conv_img, src_img.shape, anti_aliasing=True)
        score, _ = ssim(src_img, conv_img, full=True, data_range=1.0)
        scores.append(score)

    src_doc.close()
    conv_doc.close()
    mean_ssim = float(np.mean(scores))
    print(f"\nSSIM per slide: {[f'{s:.3f}' for s in scores]}")
    print(f"Mean SSIM: {mean_ssim:.3f} (threshold: {SSIM_THRESHOLD})")
    assert mean_ssim >= SSIM_THRESHOLD, (
        f"Mean SSIM {mean_ssim:.3f} below threshold {SSIM_THRESHOLD}"
    )

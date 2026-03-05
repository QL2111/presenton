"""Hybrid PDF → PPTX: screenshot background + semantic text overlay."""
import io
import os
import fitz
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor

from proto_semantic import (
    int_to_rgb,
    fitz_color_to_rgb,
    map_font_name,
    merge_text_blocks,
)

PAGE_AREA_THRESHOLD = 0.80  # skip images covering >80% of slide (already in bg)


def render_bg_erased(page, text_bboxes: list, dpi: int = 150) -> bytes:
    """Render PDF page as JPEG with text removed using PDF redaction (preserves bg images/gradients)."""
    tmp = fitz.open()
    tmp.insert_pdf(page.parent, from_page=page.number, to_page=page.number)
    cp = tmp[0]
    for bbox in text_bboxes:
        cp.add_redact_annot(fitz.Rect(bbox), fill=None)  # None = no fill, bg shows through
    cp.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE, graphics=fitz.PDF_REDACT_LINE_ART_NONE)
    pix = cp.get_pixmap(matrix=fitz.Matrix(dpi / 72, dpi / 72), colorspace=fitz.csRGB)
    tmp.close()
    buf = io.BytesIO()
    img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, 3)
    Image.fromarray(img, "RGB").save(buf, format="JPEG", quality=92)
    return buf.getvalue()


def add_images(slide, page, doc, scale_x, scale_y):
    """Add embedded sub-images on top of background, skipping full-page images."""
    page_area = page.rect.width * page.rect.height
    for img_info in page.get_images(full=True):
        xref = img_info[0]
        try:
            rects = page.get_image_rects(xref)
            if not rects:
                continue
            rect = rects[0]
            img_area = rect.width * rect.height
            if img_area > PAGE_AREA_THRESHOLD * page_area:
                continue  # full-page image already captured in background JPEG
            clipped = fitz.Rect(
                max(rect.x0, 0), max(rect.y0, 0),
                min(rect.x1, page.rect.width), min(rect.y1, page.rect.height),
            )
            if clipped.is_empty or clipped.width < 5 or clipped.height < 5:
                continue
            img_data = doc.extract_image(xref)
            if not img_data:
                continue
            left = Pt(clipped.x0 * scale_x)
            top = Pt(clipped.y0 * scale_y)
            width = Pt(clipped.width * scale_x)
            height = Pt(clipped.height * scale_y)
            slide.shapes.add_picture(io.BytesIO(img_data["image"]), left, top, width, height)
        except Exception as e:
            print(f"  [warn] image xref={xref}: {e}")


def add_text_overlay(slide, page, scale_x, scale_y):
    """Add transparent text boxes extracted from PDF over the background."""
    raw_blocks = [b for b in page.get_text("dict")["blocks"] if b.get("type") == 0]
    blocks = merge_text_blocks(raw_blocks)

    for block in blocks:
        bbox = block["bbox"]
        left   = Pt(bbox[0] * scale_x)
        top    = Pt(bbox[1] * scale_y)
        width  = Pt((bbox[2] - bbox[0]) * scale_x)
        height = Pt((bbox[3] - bbox[1]) * scale_y)
        if width < Pt(1) or height < Pt(1):
            continue

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = False
        txBox.line.fill.background()  # no border

        first_para = True
        for line in block.get("lines", []):
            para = tf.paragraphs[0] if first_para else tf.add_paragraph()
            first_para = False
            for span in line.get("spans", []):
                text = span.get("text", "")
                if not text.strip():
                    continue
                run = para.add_run()
                run.text = text
                fn = span.get("font", "")
                fn_lower = fn.lower()
                run.font.name   = map_font_name(fn)
                run.font.size   = Pt(span.get("size", 12))
                run.font.bold   = "bold" in fn_lower
                run.font.italic = "italic" in fn_lower or "oblique" in fn_lower
                c = span.get("color", 0)
                run.font.color.rgb = int_to_rgb(int(c)) if isinstance(c, (int, float)) else fitz_color_to_rgb(c)


def convert(pdf_path: str, output_path: str | None = None, dpi: int = 150) -> str:
    if not os.path.exists(pdf_path):
        raise FileNotFoundError(f"PDF not found: {pdf_path}")
    if not pdf_path.lower().endswith(".pdf"):
        raise ValueError(f"Not a PDF file: {pdf_path}")
    if output_path is None:
        output_path = os.path.splitext(pdf_path)[0] + "_hybrid.pptx"

    doc = fitz.open(pdf_path)
    prs = Presentation()
    pw = doc[0].rect.width
    ph = doc[0].rect.height
    prs.slide_width  = Pt(pw)
    prs.slide_height = Pt(ph)
    layout = prs.slide_layouts[6]  # blank

    for page_num, page in enumerate(doc):
        slide = prs.slides.add_slide(layout)
        sx = pw / page.rect.width
        sy = ph / page.rect.height

        # Collect span-level bboxes (tighter than block bboxes, preserves inter-line whitespace)
        raw_blocks = [b for b in page.get_text("dict")["blocks"] if b.get("type") == 0]
        text_bboxes = [
            span["bbox"]
            for b in raw_blocks
            for line in b.get("lines", [])
            for span in line.get("spans", [])
            if span.get("text", "").strip()
        ]

        # z=0 — background screenshot with text regions erased
        jpeg = render_bg_erased(page, text_bboxes, dpi)
        slide.shapes.add_picture(io.BytesIO(jpeg), Emu(0), Emu(0),
                                 width=prs.slide_width, height=prs.slide_height)

        add_text_overlay(slide, page, sx, sy)    # z=1 — transparent text overlay

        shapes = slide.shapes
        n_img  = sum(1 for s in shapes if s.shape_type == 13)
        n_text = sum(1 for s in shapes if s.shape_type == 17)
        print(f"Slide {page_num+1}: {len(shapes)} shapes "
              f"(1 bg + {n_img - 1} sub-img + {n_text} text)")

    doc.close()
    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    pdf_path  = os.path.join(os.path.dirname(__file__), "test_pdf_ia.pdf")
    pptx_path = os.path.join(os.path.dirname(__file__), "output", "test_output_hybrid.pptx")
    os.makedirs(os.path.dirname(pptx_path), exist_ok=True)
    out = convert(pdf_path, pptx_path)
    print(f"\nPPTX saved: {out}")
    print(f"File size: {os.path.getsize(pptx_path):,} bytes")
